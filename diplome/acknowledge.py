import smtplib
from email.message import EmailMessage
from getpass import getpass
from pptx import Presentation
import os
import comtypes.client  # Only works on Windows

# --- User input ---
name = input("Enter your name: ").strip()
email_to = input("Enter recipient's email: ").strip()
sender_email = "devmihai92@gmail.com"
app_password = "octz mybo ckmr nzsd"

# --- File paths ---
pptx_template = "diploma.pptx"
pptx_output = "diploma_filled.pptx"
pdf_output = "diploma_filled.pdf"

# --- Step 1: Replace 'NAME' in the .pptx ---
def personalize_pptx(template_path, output_path, name):
    prs = Presentation(template_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                if "NAME" in text:
                    shape.text_frame.text = text.replace("NAME", name)
    prs.save(output_path)

# --- Step 2: Convert to PDF using PowerPoint COM automation (Windows only) ---
def convert_pptx_to_pdf(pptx_path, pdf_path):
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = 1
    deck = powerpoint.Presentations.Open(os.path.abspath(pptx_path))
    deck.SaveAs(os.path.abspath(pdf_path), 32)  # 32 = pdf
    deck.Close()
    powerpoint.Quit()

# --- Step 3: Send email with PDF attached ---
def send_email_with_attachment(sender, password, recipient, subject, body, attachment_path):
    msg = EmailMessage()
    msg['Subject'] = subject
    msg['From'] = sender
    msg['To'] = recipient
    msg.set_content(body)

    with open(attachment_path, "rb") as f:
        file_data = f.read()
        msg.add_attachment(file_data, maintype='application', subtype='pdf', filename=os.path.basename(attachment_path))

    with smtplib.SMTP("smtp.gmail.com", 587) as server:
        server.starttls()
        server.login(sender, password)
        server.send_message(msg)
        print("✅ Email sent successfully.")

# --- Run the full flow ---
personalize_pptx(pptx_template, pptx_output, name)
convert_pptx_to_pdf(pptx_output, pdf_output)
send_email_with_attachment(
    sender_email,
    app_password,
    email_to,
    "Your Certificate",
    f"Hi {name},\n\nAttached is your certificate for the Agentic AI Systems module.",
    pdf_output
)
